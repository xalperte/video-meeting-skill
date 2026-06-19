import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "scripts"))

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAVE_PPTX = True
except ImportError:
    HAVE_PPTX = False


def _record():
    return {
        "meeting": {"title": "Sprint", "date": "2026-06-18", "type": "grooming"},
        "summary": {"tldr": "All good", "sections": [
            {"category": "Notes", "points": ["a", "b"]}]},
        "action_items": [{"title": "Do X", "type": "explicit", "assignee": "Al"}],
    }


@unittest.skipUnless(HAVE_PPTX, "python-pptx required")
class SlidesHardening(unittest.TestCase):
    def _write_record(self, d):
        import json
        rec = os.path.join(d, "rec.json")
        with open(rec, "w") as fh:
            json.dump(_record(), fh)
        return rec

    def test_clears_preexisting_slides(self):
        from render_slides import render
        with tempfile.TemporaryDirectory() as d:
            # base deck seeded with 3 example slides
            tdir = os.path.join(d, "tpl"); os.mkdir(tdir)
            base = Presentation()
            for _ in range(3):
                base.slides.add_slide(base.slide_layouts[6])  # blank
            base.save(os.path.join(tdir, "slides.pptx"))
            out = os.path.join(d, "out.pptx")
            render(self._write_record(d), out, ["pptx"], "soffice", tdir)
            prs = Presentation(out)
            # only generated slides remain (title + overview + notes + 2 action groups = 5)
            # crucially: NOT 3 + generated
            titles = []
            for s in prs.slides:
                if s.shapes.title and s.shapes.title.text:
                    titles.append(s.shapes.title.text)
            self.assertIn("Sprint", titles)
            self.assertLessEqual(len(prs.slides._sldIdLst), 6)
            self.assertGreaterEqual(len(prs.slides._sldIdLst), 4)

    def test_layout_without_placeholders_does_not_crash(self):
        from render_slides import render
        with tempfile.TemporaryDirectory() as d:
            # base deck whose only usable layout is blank (no title/body placeholders)
            tdir = os.path.join(d, "tpl"); os.mkdir(tdir)
            base = Presentation()
            # template.yaml points both layouts at "Blank" (idx 6, no placeholders)
            with open(os.path.join(tdir, "template.yaml"), "w") as fh:
                fh.write('slides:\n  title_layout: "Blank"\n  content_layout: "Blank"\n')
            base.save(os.path.join(tdir, "slides.pptx"))
            out = os.path.join(d, "out.pptx")
            render(self._write_record(d), out, ["pptx"], "soffice", tdir)  # must not raise
            prs = Presentation(out)
            self.assertGreater(len(prs.slides._sldIdLst), 0)
            # the title text was placed (in a textbox fallback)
            alltext = []
            for s in prs.slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        alltext.append(sh.text_frame.text)
            self.assertTrue(any("Sprint" in t for t in alltext))


if __name__ == "__main__":
    unittest.main()
