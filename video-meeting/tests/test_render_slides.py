import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "scripts"))

try:
    from pptx import Presentation
    HAVE_PPTX = True
except ImportError:
    HAVE_PPTX = False


@unittest.skipUnless(HAVE_PPTX, "python-pptx required")
class SlidesTemplate(unittest.TestCase):
    def _record(self):
        return {
            "meeting": {"title": "Sprint", "date": "2026-06-18", "type": "grooming"},
            "summary": {"tldr": "All good", "sections": [
                {"category": "Notes", "points": ["a", "b"]}]},
            "action_items": [{"title": "Do X", "type": "explicit", "assignee": "Al"}],
        }

    def _make_base_template(self, d):
        """A template dir whose slides.pptx is the default python-pptx deck
        (which DOES contain 'Title Slide' and 'Title and Content' layouts)."""
        tdir = os.path.join(d, "tpl")
        os.mkdir(tdir)
        Presentation().save(os.path.join(tdir, "slides.pptx"))
        return tdir

    def test_renders_from_base_template(self):
        import json
        from render_slides import render
        with tempfile.TemporaryDirectory() as d:
            tdir = self._make_base_template(d)
            rec = os.path.join(d, "rec.json")
            with open(rec, "w") as fh:
                json.dump(self._record(), fh)
            out = os.path.join(d, "slides.pptx")
            render(rec, out, ["pptx"], "soffice", tdir)
            self.assertTrue(os.path.isfile(out))
            self.assertGreater(len(Presentation(out).slides._sldIdLst), 0)

    def test_no_template_dir_uses_default_deck(self):
        import json
        from render_slides import render
        with tempfile.TemporaryDirectory() as d:
            rec = os.path.join(d, "rec.json")
            with open(rec, "w") as fh:
                json.dump(self._record(), fh)
            out = os.path.join(d, "slides.pptx")
            render(rec, out, ["pptx"], "soffice", None)
            self.assertTrue(os.path.isfile(out))


if __name__ == "__main__":
    unittest.main()
