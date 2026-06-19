#!/usr/bin/env python3
"""Generate the starter presentation templates (run in render-env).

Produces clean, editable base files for each starter palette:
  templates/presentation/<name>/slides.pptx
  templates/presentation/<name>/report.docx
  templates/presentation/<name>/template.yaml
  templates/presentation/<name>/logo.png   (1x1 transparent placeholder)

Re-runnable and idempotent. Output is committed so a bare checkout has working
templates. Replace slides.pptx/report.docx with a corporate file to rebrand.

Usage:
  $RENDER_PY build_starter_templates.py [--root <skill_root>]
"""
import argparse
import base64
import copy
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.util import Pt as PPt
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches as PInches
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGB
from docx.enum.style import WD_STYLE_TYPE  # noqa: F401 (kept for clarity)
from docx.oxml.ns import qn

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from render_report import clear_body  # reuse body-clearing logic

HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_ROOT = os.path.normpath(os.path.join(HERE, ".."))

# 1x1 transparent PNG (placeholder logo).
_LOGO_PNG = base64.b64decode(
    b"iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
    b"+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg==")

STARTERS = {
    "executive": {"name": "Executive", "accent": "1A2238", "accent2": "B08D57",
                  "desc": "Navy + gold, Cambria headings — corporate executive brand",
                  "heading_font": "Cambria", "body_font": "Calibri"},
    "client": {"name": "Client", "accent": "117A65",
               "desc": "Teal, sans — client-facing, logo-forward",
               "heading_font": "Calibri", "body_font": "Calibri"},
    "internal": {"name": "Internal", "accent": "595959",
                 "desc": "Minimal grey — internal team default",
                 "heading_font": "Calibri", "body_font": "Calibri"},
}


def _style_slide_master(prs, cfg):
    """Color the title placeholders on each layout with the accent."""
    accent = PRGB.from_string(cfg["accent"])
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 0 and ph.has_text_frame:  # title
                for para in ph.text_frame.paragraphs:
                    para.font.color.rgb = accent
                    para.font.name = cfg["heading_font"]
                    para.font.size = PPt(32)


def _theme_part(prs):
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    master_part = prs.slide_masters[0].part
    for rel in master_part.rels.values():
        if rel.reltype == RT.THEME:
            return rel.target_part
    return None


def _brand_theme(prs, cfg):
    """Best-effort: set theme dk2/accent1 colors and major font. Falls back
    silently (keeps default theme) if the structure is unexpected."""
    part = _theme_part(prs)
    if part is None:
        return
    try:
        root = part._element
        a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        clr = root.find(f"{a}themeElements/{a}clrScheme")
        for tag, hexv in (("dk2", cfg["accent"]), ("accent1", cfg["accent"]),
                          ("accent2", cfg.get("accent2", cfg["accent"]))):
            node = clr.find(f"{a}{tag}/{a}srgbClr")
            if node is not None:
                node.set("val", hexv)
        fonts = root.find(f"{a}themeElements/{a}fontScheme")
        for slot in ("majorFont", "minorFont"):
            latin = fonts.find(f"{a}{slot}/{a}latin")
            if latin is not None and slot == "majorFont":
                latin.set("typeface", cfg["heading_font"])
    except Exception as exc:  # noqa: BLE001 - best effort, never abort the build
        sys.stderr.write(f"  warn: theme recolor skipped: {exc}\n")


def _add_master_band_and_logo(prs, cfg, logo_path):
    master = prs.slide_masters[0]
    accent = PRGB.from_string(cfg["accent"])
    try:
        band = master.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PInches(0.35))
        band.fill.solid(); band.fill.fore_color.rgb = accent
        band.line.fill.background()
    except Exception as exc:  # noqa: BLE001
        sys.stderr.write(f"  warn: title band skipped: {exc}\n")
    if logo_path and os.path.isfile(logo_path):
        try:
            master.shapes.add_picture(
                logo_path, prs.slide_width - PInches(1.4), PInches(0.05),
                height=PInches(0.5))
        except Exception as exc:  # noqa: BLE001
            sys.stderr.write(f"  warn: logo placement skipped: {exc}\n")


def build_slides(folder, cfg):
    prs = Presentation()
    _brand_theme(prs, cfg)
    _style_slide_master(prs, cfg)  # guaranteed title color/font on layouts
    _add_master_band_and_logo(prs, cfg, os.path.join(folder, "logo.png"))
    prs.save(os.path.join(folder, "slides.pptx"))


def _ensure_styles(doc, names):
    """Copy missing builtin styles (by w:name val) from a blank Document."""
    have = {s.name for s in doc.styles}
    missing = [n for n in names if n not in have]
    if not missing:
        return
    ref = Document()
    ref_styles = ref.styles.element
    doc_styles = doc.styles.element
    for st in ref_styles.findall(qn("w:style")):
        n = st.find(qn("w:name"))
        if n is not None and n.get(qn("w:val")) in missing:
            doc_styles.append(copy.deepcopy(st))


def build_report(folder, cfg):
    source = os.path.join(folder, "source", "Executive_Template.docx")
    if os.path.isfile(source):
        doc = Document(source)
        clear_body(doc)                       # keep header/footer + section setup
        _ensure_styles(doc, ["Table Grid", "List Bullet"])
    else:
        doc = Document()
        if "Normal" in {s.name for s in doc.styles}:
            normal = doc.styles["Normal"]
            normal.font.name, normal.font.size = cfg["body_font"], DPt(11)

    accent = DRGB.from_string(cfg["accent"])
    target_levels = {"Title", "Heading 1", "Heading 2"}
    for st in doc.styles:
        if st.name in target_levels:
            st.font.name = cfg["heading_font"]
            st.font.color.rgb = accent
    doc.save(os.path.join(folder, "report.docx"))


def build_template_yaml(folder, cfg):
    accent = cfg["accent"]
    accent2 = cfg.get("accent2", accent)
    text = (
        f'name: "{cfg["name"]}"\n'
        f'description: "{cfg["desc"]}"\n'
        f'accent: "{accent2}"\n'
        "slides:\n"
        '  title_layout: "Title Slide"\n'
        '  content_layout: "Title and Content"\n'
        "xlsx:\n"
        f'  header_fill: "{accent}"\n'
        '  header_font_color: "FFFFFF"\n'
        '  explicit_fill: "E2EFDA"\n'
        '  suggested_fill: "FFF2CC"\n'
        f'  font_name: "{cfg["body_font"]}"\n'
        "  banner: false\n"
    )
    with open(os.path.join(folder, "template.yaml"), "w", encoding="utf-8") as fh:
        fh.write(text)


def build_all(root):
    base = os.path.join(root, "templates", "presentation")
    for key, cfg in STARTERS.items():
        folder = os.path.join(base, key)
        os.makedirs(folder, exist_ok=True)
        logo = os.path.join(folder, "logo.png")
        if not os.path.isfile(logo):
            with open(logo, "wb") as fh:
                fh.write(_LOGO_PNG)
        build_slides(folder, cfg)
        build_report(folder, cfg)
        build_template_yaml(folder, cfg)
        sys.stderr.write(f"  built template: {key} -> {folder}\n")


def main():
    ap = argparse.ArgumentParser(description="Generate starter presentation templates.")
    ap.add_argument("--root", default=DEFAULT_ROOT)
    args = ap.parse_args()
    build_all(args.root)


if __name__ == "__main__":
    main()
