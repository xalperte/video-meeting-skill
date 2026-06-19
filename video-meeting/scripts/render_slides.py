#!/usr/bin/env python3
"""
Renderer — slides.pptx (+ optional .odp) from meeting_record.json.

Runs in render-env (python-pptx). The deck is a presentation VIEW of the record:
title -> overview (TL;DR) -> categorized summary -> action items split into two
groups (explicit vs AI-suggested). The .odp copy is produced by converting the
.pptx with LibreOffice headless.

Usage:
  $RENDER_PY render_slides.py --record meeting_record.json \
     --out-pptx slides.pptx --formats pptx odp --libreoffice /usr/bin/soffice
"""
import argparse
import json
import os
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from templates import find_layout_by_name, load_template_yaml, slides_layouts  # noqa: E402

MAX_BULLETS = 7  # split a section across slides beyond this many points


def clear_slides(prs):
    """Remove all existing slides from prs, keeping masters/layouts/theme."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rId = sld_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        sld_id_lst.remove(sld_id)


def _title_placeholder(slide):
    return slide.shapes.title  # may be None


def _body_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    return None


def _set_title(slide, title):
    ph = _title_placeholder(slide)
    if ph is not None:
        ph.text = title
        return
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    box.text_frame.text = title


def _body_frame(slide):
    ph = _body_placeholder(slide)
    if ph is not None and ph.has_text_frame:
        return ph.text_frame
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    return box.text_frame


def add_title_slide(prs, layout, title, subtitle):
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    sub = _body_placeholder(slide)
    if sub is not None:
        sub.text = subtitle
    elif subtitle:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(1))
        box.text_frame.text = subtitle


def add_bullets_slide(prs, layout, title, bullets):
    """bullets: list of (text, level)."""
    slide = prs.slides.add_slide(layout)
    _set_title(slide, title)
    body = _body_frame(slide)
    body.clear()
    if not bullets:
        body.paragraphs[0].text = "(none)"
        return
    for i, (text, level) in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = min(level, 4)
        p.font.size = Pt(18)


def chunked(title, bullets):
    """Yield (slide_title, bullets) splitting long lists across slides."""
    if len(bullets) <= MAX_BULLETS:
        yield title, bullets
        return
    n = 0
    for i in range(0, len(bullets), MAX_BULLETS):
        n += 1
        suffix = "" if i == 0 else f" (cont. {n})"
        yield title + suffix, bullets[i:i + MAX_BULLETS]


def item_line(t):
    who = t.get("assignee") or "unassigned"
    pri = t.get("priority", "")
    return f"{t.get('title','').strip()} — {who}" + (f"  [{pri}]" if pri else "")


def build(prs, record, layout_cfg):
    title_layout, _ = find_layout_by_name(
        prs.slide_layouts, layout_cfg["title_layout"], 0)
    content_layout, _ = find_layout_by_name(
        prs.slide_layouts, layout_cfg["content_layout"], 1)

    m = record.get("meeting", {})
    title = m.get("title", "Meeting")
    sub_bits = [b for b in (m.get("date"), (m.get("type") or "").title()) if b]
    add_title_slide(prs, title_layout, title, "  ·  ".join(sub_bits))

    tldr = record.get("summary", {}).get("tldr", "").strip()
    if tldr:
        add_bullets_slide(prs, content_layout, "Overview", [(tldr, 0)])

    for sec in record.get("summary", {}).get("sections", []):
        cat = sec.get("category", "Notes")
        bullets = [(p, 0) for p in sec.get("points", [])]
        for stitle, sb in chunked(cat, bullets):
            add_bullets_slide(prs, content_layout, stitle, sb)

    items = record.get("action_items", [])
    explicit = [item_line(t) for t in items if t.get("type") == "explicit"]
    suggested = [item_line(t) for t in items if t.get("type") == "ai_suggested"]
    for stitle, sb in chunked("Action Items", [(x, 0) for x in explicit]):
        add_bullets_slide(prs, content_layout, stitle, sb or [("(none)", 0)])
    for stitle, sb in chunked("Suggested (AI)", [(x, 0) for x in suggested]):
        add_bullets_slide(prs, content_layout, stitle, sb or [("(none)", 0)])


def convert_to(soffice, fmt, src, outdir):
    cmd = [soffice, "--headless", "--convert-to", fmt, "--outdir", outdir, src]
    p = subprocess.run(cmd, capture_output=True, text=True, check=False)
    if p.returncode != 0:
        sys.stderr.write(f"  warn: {fmt} conversion failed: {p.stderr.strip()}\n")
        return None
    out = os.path.join(outdir, os.path.splitext(os.path.basename(src))[0] + f".{fmt}")
    return out if os.path.isfile(out) else None


def render(record_path, out_pptx, formats, libreoffice, template_dir=None):
    with open(record_path, "r", encoding="utf-8") as fh:
        record = json.load(fh)

    data = load_template_yaml(template_dir) if template_dir else {}
    base = os.path.join(template_dir, "slides.pptx") if template_dir else None
    used_base = bool(base and os.path.isfile(base))
    prs = Presentation(base) if used_base else Presentation()
    if used_base:
        clear_slides(prs)

    build(prs, record, slides_layouts(data))
    os.makedirs(os.path.dirname(os.path.abspath(out_pptx)) or ".", exist_ok=True)
    prs.save(out_pptx)
    sys.stderr.write(f"  slides.pptx: {len(prs.slides._sldIdLst)} slides -> {out_pptx}\n")

    if "odp" in formats:
        outdir = os.path.dirname(os.path.abspath(out_pptx)) or "."
        odp = convert_to(libreoffice, "odp", out_pptx, outdir)
        if odp:
            sys.stderr.write(f"  slides.odp -> {odp}\n")


def main():
    ap = argparse.ArgumentParser(description="Render slides from the meeting record.")
    ap.add_argument("--record", required=True)
    ap.add_argument("--out-pptx", required=True)
    ap.add_argument("--formats", nargs="*", default=["pptx", "odp"])
    ap.add_argument("--libreoffice", default="soffice")
    ap.add_argument("--template-dir", default=None)
    args = ap.parse_args()

    if not os.path.isfile(args.record):
        sys.exit(f"record not found: {args.record}")
    render(args.record, args.out_pptx, args.formats, args.libreoffice, args.template_dir)


if __name__ == "__main__":
    main()
